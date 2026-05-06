#!/usr/bin/env python3
"""Generate a Security Query Hub – Data Flow Diagram PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# ── Palette ──────────────────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x11, 0x17)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0x8B, 0x95, 0xA5)
ACCENT    = RGBColor(0x58, 0xA6, 0xFF)
GREEN     = RGBColor(0x3F, 0xB9, 0x50)
ORANGE    = RGBColor(0xF0, 0x88, 0x3E)
RED       = RGBColor(0xF8, 0x53, 0x49)
PURPLE    = RGBColor(0xBC, 0x8C, 0xFF)
TEAL      = RGBColor(0x39, 0xD3, 0xC8)
CARD_BG   = RGBColor(0x16, 0x1B, 0x22)
CARD_BDR  = RGBColor(0x30, 0x36, 0x3D)
ARROW_CLR = RGBColor(0x48, 0x8B, 0xD4)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def _set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_box(slide, left, top, w, h, fill=CARD_BG, border=CARD_BDR, text="",
             font_size=10, font_color=WHITE, bold=False, align=PP_ALIGN.LEFT, line_w=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = line_w
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    if text:
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = align
    return shape

def _add_text(slide, left, top, w, h, text, size=12, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def _add_para(tf, text, size=10, color=WHITE, bold=False, space_before=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = space_before
    return p

def _arrow(slide, x1, y1, x2, y2, color=ARROW_CLR, width=Pt(2)):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    conn.line.color.rgb = color
    conn.line.width = width
    # end arrow
    conn.begin_x = x1
    conn.begin_y = y1
    conn.end_x = x2
    conn.end_y = y2
    return conn


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_set_bg(sl)
_add_text(sl, Inches(1), Inches(2.2), Inches(11), Inches(1),
          "Security Query Hub", size=44, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
_add_text(sl, Inches(1), Inches(3.4), Inches(11), Inches(0.6),
          "Data Flow Architecture", size=24, color=WHITE, align=PP_ALIGN.CENTER)
_add_text(sl, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
          "FastAPI  |  SQLite/SQLCipher  |  SentinelOne Deep Visibility  |  Vanilla JS SPA",
          size=14, color=GREY, align=PP_ALIGN.CENTER)
_add_text(sl, Inches(1), Inches(5.8), Inches(11), Inches(0.4),
          "Warner Bros. Discovery – Security Engineering", size=12, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 – AWS Infrastructure Overview
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(sl)
_add_text(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
          "AWS Infrastructure Overview", size=28, color=ACCENT, bold=True)

# ── Left column: Network path ──
# User / VPN box
_add_box(sl, Inches(0.4), Inches(1.3), Inches(2.6), Inches(1.4), border=ACCENT,
         text="", font_size=10)
user_box = sl.shapes[-1]
utf = user_box.text_frame
utf.paragraphs[0].text = "Security Analysts"
utf.paragraphs[0].font.size = Pt(13)
utf.paragraphs[0].font.color.rgb = ACCENT
utf.paragraphs[0].font.bold = True
_add_para(utf, "Corporate laptops", size=9, color=GREY, space_before=Pt(4))
_add_para(utf, "Connected via VPN", size=9, color=GREY, space_before=Pt(2))
_add_para(utf, "HTTPS to port 80", size=9, color=GREY, space_before=Pt(2))

_add_text(sl, Inches(3.05), Inches(1.7), Inches(0.4), Inches(0.4), "→", size=20, color=ARROW_CLR, bold=True)

# Security Group box
_add_box(sl, Inches(3.5), Inches(1.3), Inches(2.6), Inches(1.4), border=ORANGE,
         text="", font_size=10)
sg_box = sl.shapes[-1]
sgtf = sg_box.text_frame
sgtf.paragraphs[0].text = "Security Group / ACL"
sgtf.paragraphs[0].font.size = Pt(13)
sgtf.paragraphs[0].font.color.rgb = ORANGE
sgtf.paragraphs[0].font.bold = True
_add_para(sgtf, "Inbound: port 80 from", size=9, color=GREY, space_before=Pt(4))
_add_para(sgtf, "  VPN subnets only", size=9, color=GREY, space_before=Pt(2))
_add_para(sgtf, "Outbound: S1 API (443)", size=9, color=GREY, space_before=Pt(2))

_add_text(sl, Inches(6.15), Inches(1.7), Inches(0.4), Inches(0.4), "→", size=20, color=ARROW_CLR, bold=True)

# ── EC2 instance (large box) ──
ec2 = _add_box(sl, Inches(6.6), Inches(1.3), Inches(6.3), Inches(5.5), border=GREEN,
               text="", font_size=10)
ec2tf = ec2.text_frame
ec2tf.paragraphs[0].text = "AWS EC2 Instance"
ec2tf.paragraphs[0].font.size = Pt(16)
ec2tf.paragraphs[0].font.color.rgb = GREEN
ec2tf.paragraphs[0].font.bold = True
_add_para(ec2tf, "Ubuntu 24.04 LTS  |  t3.medium (2 vCPU, 4 GB RAM)  |  IP: 3.145.46.91", size=10, color=GREY, space_before=Pt(4))

# Inner boxes inside EC2
# -- systemd service --
_add_box(sl, Inches(6.9), Inches(2.6), Inches(2.7), Inches(1.8), border=TEAL)
svc = sl.shapes[-1]
svctf = svc.text_frame
svctf.paragraphs[0].text = "systemd: sqh.service"
svctf.paragraphs[0].font.size = Pt(12)
svctf.paragraphs[0].font.color.rgb = TEAL
svctf.paragraphs[0].font.bold = True
_add_para(svctf, "User: sqh (no-login)", size=9, color=GREY, space_before=Pt(4))
_add_para(svctf, "ExecStart: venv/bin/python run.py", size=8, color=GREY, space_before=Pt(2))
_add_para(svctf, "Restart=always, RestartSec=5", size=8, color=GREY, space_before=Pt(2))
_add_para(svctf, "ProtectSystem=strict", size=8, color=GREY, space_before=Pt(2))
_add_para(svctf, "NoNewPrivileges=yes", size=8, color=GREY, space_before=Pt(2))

# -- Uvicorn + FastAPI --
_add_box(sl, Inches(9.9), Inches(2.6), Inches(2.7), Inches(1.8), border=ACCENT)
uvi = sl.shapes[-1]
uvitf = uvi.text_frame
uvitf.paragraphs[0].text = "Uvicorn + FastAPI"
uvitf.paragraphs[0].font.size = Pt(12)
uvitf.paragraphs[0].font.color.rgb = ACCENT
uvitf.paragraphs[0].font.bold = True
_add_para(uvitf, "Listens on 0.0.0.0:80", size=9, color=GREY, space_before=Pt(4))
_add_para(uvitf, "ASGI async server", size=9, color=GREY, space_before=Pt(2))
_add_para(uvitf, "Serves API + static files", size=9, color=GREY, space_before=Pt(2))
_add_para(uvitf, "CAP_NET_BIND_SERVICE", size=8, color=GREY, space_before=Pt(2))
_add_para(uvitf, "  (bind port 80 without root)", size=8, color=GREY, space_before=Pt(2))

# -- EBS volume --
_add_box(sl, Inches(6.9), Inches(4.8), Inches(2.7), Inches(1.7), border=PURPLE)
ebs = sl.shapes[-1]
ebstf = ebs.text_frame
ebstf.paragraphs[0].text = "EBS gp3 Volume (100 GB)"
ebstf.paragraphs[0].font.size = Pt(12)
ebstf.paragraphs[0].font.color.rgb = PURPLE
ebstf.paragraphs[0].font.bold = True
_add_para(ebstf, "Encrypted at rest", size=9, color=GREY, space_before=Pt(4))
_add_para(ebstf, "/opt/sqh/data/sqh.db", size=8, color=GREY, space_before=Pt(2))
_add_para(ebstf, "  SQLCipher encrypted DB", size=8, color=GREY, space_before=Pt(2))
_add_para(ebstf, "/opt/sqh/logs/", size=8, color=GREY, space_before=Pt(2))

# -- Secrets --
_add_box(sl, Inches(9.9), Inches(4.8), Inches(2.7), Inches(1.7), border=RED)
sec = sl.shapes[-1]
sectf = sec.text_frame
sectf.paragraphs[0].text = "Secrets  (/etc/sqh/env)"
sectf.paragraphs[0].font.size = Pt(12)
sectf.paragraphs[0].font.color.rgb = RED
sectf.paragraphs[0].font.bold = True
_add_para(sectf, "SQH_ENCRYPTION_KEY (Fernet)", size=9, color=GREY, space_before=Pt(4))
_add_para(sectf, "SQH_DB_PASSPHRASE (SQLCipher)", size=9, color=GREY, space_before=Pt(2))
_add_para(sectf, "root:sqh 640 permissions", size=9, color=GREY, space_before=Pt(2))
_add_para(sectf, "Loaded via EnvironmentFile=", size=8, color=GREY, space_before=Pt(2))

# ── Bottom left: S1 API ──
_add_box(sl, Inches(0.4), Inches(3.2), Inches(2.6), Inches(1.4), border=ORANGE)
s1 = sl.shapes[-1]
s1tf = s1.text_frame
s1tf.paragraphs[0].text = "SentinelOne Cloud"
s1tf.paragraphs[0].font.size = Pt(13)
s1tf.paragraphs[0].font.color.rgb = ORANGE
s1tf.paragraphs[0].font.bold = True
_add_para(s1tf, "Deep Visibility API", size=9, color=GREY, space_before=Pt(4))
_add_para(s1tf, "HTTPS (port 443)", size=9, color=GREY, space_before=Pt(2))
_add_para(s1tf, "API token authentication", size=9, color=GREY, space_before=Pt(2))

_add_text(sl, Inches(3.05), Inches(3.5), Inches(3.5), Inches(0.4), "← EC2 outbound (443) →", size=10, color=ARROW_CLR, bold=True)

# ── Bottom left: Directory layout ──
_add_box(sl, Inches(0.4), Inches(5.1), Inches(5.7), Inches(1.7), border=CARD_BDR)
dirs = sl.shapes[-1]
dirtf = dirs.text_frame
dirtf.paragraphs[0].text = "Server Directory Layout"
dirtf.paragraphs[0].font.size = Pt(12)
dirtf.paragraphs[0].font.color.rgb = TEAL
dirtf.paragraphs[0].font.bold = True
_add_para(dirtf, "/opt/sqh/                Application root (owned by sqh user)", size=9, color=GREY, space_before=Pt(4))
_add_para(dirtf, "/opt/sqh/app/            FastAPI application code + static assets", size=9, color=GREY, space_before=Pt(2))
_add_para(dirtf, "/opt/sqh/venv/           Python virtual environment", size=9, color=GREY, space_before=Pt(2))
_add_para(dirtf, "/opt/sqh/data/sqh.db     SQLCipher encrypted database", size=9, color=GREY, space_before=Pt(2))
_add_para(dirtf, "/opt/sqh/logs/           Rotating log files (10MB x 5)", size=9, color=GREY, space_before=Pt(2))
_add_para(dirtf, "/etc/sqh/env             Secrets (Fernet key, DB passphrase)", size=9, color=GREY, space_before=Pt(2))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 – High-Level Software Architecture
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(sl)
_add_text(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
          "High-Level Architecture", size=28, color=ACCENT, bold=True)

# --- Browser box ---
bx = _add_box(sl, Inches(0.5), Inches(1.2), Inches(2.8), Inches(2.2), border=ACCENT)
tf = bx.text_frame
tf.paragraphs[0].text = "Browser (SPA)"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = ACCENT
tf.paragraphs[0].font.bold = True
for item in ["Login / Auth UI", "Query Selector + Params", "Results Table", "AI Detection Dashboard", "Admin Panel", "History + Export"]:
    _add_para(tf, f"•  {item}", size=9, color=GREY)

# --- FastAPI box ---
bx2 = _add_box(sl, Inches(4.2), Inches(1.2), Inches(3.5), Inches(2.2), border=GREEN)
tf2 = bx2.text_frame
tf2.paragraphs[0].text = "FastAPI Backend"
tf2.paragraphs[0].font.size = Pt(14)
tf2.paragraphs[0].font.color.rgb = GREEN
tf2.paragraphs[0].font.bold = True
for item in ["/api/auth/*  (sessions, bcrypt)", "/api/queries/*  (run, cancel)", "/api/history/*  (results, export)",
             "/api/admin/*  (users, queries)", "/api/system/*  (config, disk)", "/api/ai-tools  (CRUD)"]:
    _add_para(tf2, f"•  {item}", size=9, color=GREY)

# --- S1 box ---
bx3 = _add_box(sl, Inches(8.5), Inches(1.2), Inches(3.5), Inches(2.2), border=ORANGE)
tf3 = bx3.text_frame
tf3.paragraphs[0].text = "SentinelOne API"
tf3.paragraphs[0].font.size = Pt(14)
tf3.paragraphs[0].font.color.rgb = ORANGE
tf3.paragraphs[0].font.bold = True
for item in ["POST /dv/init-query", "GET  /dv/query-status", "GET  /dv/events (cursor)", "POST /dv/query-pq (PowerQuery)"]:
    _add_para(tf3, f"•  {item}", size=9, color=GREY)

# --- DB box ---
bx4 = _add_box(sl, Inches(4.2), Inches(4.0), Inches(3.5), Inches(2.8), border=PURPLE)
tf4 = bx4.text_frame
tf4.paragraphs[0].text = "SQLite + SQLCipher"
tf4.paragraphs[0].font.size = Pt(14)
tf4.paragraphs[0].font.color.rgb = PURPLE
tf4.paragraphs[0].font.bold = True
for item in ["users  /  sessions", "app_config  (encrypted keys)", "stored_queries  /  query_params",
             "query_folders", "query_history", "query_results  (gzip compressed)", "ai_tools"]:
    _add_para(tf4, f"•  {item}", size=9, color=GREY)

# --- Backend services box ---
bx5 = _add_box(sl, Inches(8.5), Inches(4.0), Inches(3.5), Inches(2.8), border=TEAL)
tf5 = bx5.text_frame
tf5.paragraphs[0].text = "Backend Services"
tf5.paragraphs[0].font.size = Pt(14)
tf5.paragraphs[0].font.color.rgb = TEAL
tf5.paragraphs[0].font.bold = True
for item in ["s1_client.py – time-slicing, polling", "fast_json.py – orjson (10x faster)",
             "auth.py – bcrypt, session mgmt", "secrets_manager.py – Fernet encrypt",
             "disk_monitor.py – FIFO cleanup", "database.py – schema, migrations"]:
    _add_para(tf5, f"•  {item}", size=9, color=GREY)

# --- Arrows (conceptual labels) ---
_add_text(sl, Inches(3.3), Inches(1.9), Inches(1), Inches(0.3), "HTTP →", size=10, color=ARROW_CLR, bold=True)
_add_text(sl, Inches(7.7), Inches(1.9), Inches(1), Inches(0.3), "API →", size=10, color=ARROW_CLR, bold=True)
_add_text(sl, Inches(7.7), Inches(4.8), Inches(1), Inches(0.3), "← uses", size=10, color=ARROW_CLR, bold=True)
_add_text(sl, Inches(5.4), Inches(3.5), Inches(1), Inches(0.3), "↕ R/W", size=10, color=ARROW_CLR, bold=True)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Query Execution Flow
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(sl)
_add_text(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
          "Query Execution Flow", size=28, color=ACCENT, bold=True)

steps = [
    ("1", "User Executes Query", "Browser sends POST /queries/{id}/run\nwith params, from_date, to_date", ACCENT),
    ("2", "Backend Queues Task", "Inserts query_history (status=running)\nSpawns asyncio background task", GREEN),
    ("3", "Time-Slice & Fan Out", "Splits date range into 8 windows\nRuns all 8 concurrently via shared\nhttpx client + semaphore", ORANGE),
    ("4", "S1 Per-Slice Flow", "init-query → poll status (0.5-2s)\n→ fetch events (cursor, 1000/page)\nAll 8 slices run in parallel", RED),
    ("5", "Merge & Dedup", "Combine all slice results\nDedup by event ID across slices\nAI queries: extra dedup by\nendpoint+process+user", PURPLE),
    ("6", "Compress & Store", "orjson serialize → gzip compress\nINSERT into query_results\nUpdate history status=success", TEAL),
]

for i, (num, title, desc, color) in enumerate(steps):
    col = i % 3
    row = i // 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.2 + row * 3.0)

    # number circle
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.45), Inches(0.45))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    ctf = circ.text_frame
    ctf.paragraphs[0].text = num
    ctf.paragraphs[0].font.size = Pt(16)
    ctf.paragraphs[0].font.color.rgb = WHITE
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # card
    card = _add_box(sl, left + Inches(0.6), top, Inches(3.4), Inches(2.4))
    ctf2 = card.text_frame
    ctf2.paragraphs[0].text = title
    ctf2.paragraphs[0].font.size = Pt(13)
    ctf2.paragraphs[0].font.color.rgb = color
    ctf2.paragraphs[0].font.bold = True
    for line in desc.split("\n"):
        _add_para(ctf2, line, size=10, color=GREY, space_before=Pt(4))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Result Loading Flow
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(sl)
_add_text(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
          "Result Loading & Display Flow", size=28, color=ACCENT, bold=True)

load_steps = [
    ("Browser polls status", "GET /queries/status/{id} every 3s\nUntil status = success", ACCENT),
    ("Fetch results", "GET /history/{id}/results?limit=100000\nFor AI dashboard, loads all events\nin a single request", GREEN),
    ("Server decompresses", "Read gzip blob from SQLite\ngzip.decompress → orjson.loads\nSlice page → orjson.dumps", ORANGE),
    ("Gzip response", "Response compressed with gzip\nfor fast network transfer\nAccept-Encoding: gzip", PURPLE),
    ("Render UI", "Regular query → sortable table\nAI query → dashboard with cards,\nChart.js doughnut, drill-down tables", TEAL),
]

for i, (title, desc, color) in enumerate(load_steps):
    left = Inches(0.4 + i * 2.55)
    top = Inches(1.5)

    card = _add_box(sl, left, top, Inches(2.35), Inches(3.0))
    ctf = card.text_frame
    # step number
    p0 = ctf.paragraphs[0]
    p0.text = f"Step {i+1}"
    p0.font.size = Pt(9)
    p0.font.color.rgb = GREY
    p0.font.bold = True
    # title
    _add_para(ctf, title, size=12, color=color, bold=True, space_before=Pt(6))
    # desc
    for line in desc.split("\n"):
        _add_para(ctf, line, size=9, color=GREY, space_before=Pt(4))

    # arrow between cards
    if i < len(load_steps) - 1:
        _add_text(sl, left + Inches(2.35), Inches(2.7), Inches(0.3), Inches(0.3),
                  "→", size=16, color=ARROW_CLR, bold=True)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Performance Optimizations
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(sl)
_add_text(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
          "Performance Optimizations", size=28, color=ACCENT, bold=True)

opts = [
    ("Concurrent Time-Slicing", "8 slices, 8 concurrent\nAll run in parallel via asyncio\nShared httpx connection pool",
     "Before: 80-120s  →  After: 15-30s", GREEN),
    ("orjson Serialization", "C-extension JSON library\n~10x faster than stdlib json\nUsed for all result I/O",
     "20K events: 3s → 0.3s", ORANGE),
    ("Gzip Compressed Storage", "Results stored gzip-compressed\nin SQLite (compresslevel=1)\n~10x smaller on disk",
     "50MB raw → ~5MB compressed", PURPLE),
    ("Fast Polling", "0.5s intervals for first 6 polls\n1s for polls 7-20, then 2s\nCatches fast queries sooner",
     "Saves 5-10s per slice", TEAL),
    ("S1QL In Contains", "Single operator for all keywords\nvs 46 separate OR clauses\nStays under S1 parse limits",
     "Query length: 2KB → 400 bytes", ACCENT),
    ("Gzip HTTP Responses", "End-to-end compression\nServer gzip → browser decompress\nFast even on slow networks",
     "Transfer: 50MB → ~5MB", RED),
]

for i, (title, desc, perf, color) in enumerate(opts):
    col = i % 3
    row = i // 3
    left = Inches(0.4 + col * 4.2)
    top = Inches(1.2 + row * 3.1)

    card = _add_box(sl, left, top, Inches(3.9), Inches(2.6))
    ctf = card.text_frame
    ctf.paragraphs[0].text = title
    ctf.paragraphs[0].font.size = Pt(14)
    ctf.paragraphs[0].font.color.rgb = color
    ctf.paragraphs[0].font.bold = True
    for line in desc.split("\n"):
        _add_para(ctf, line, size=10, color=GREY, space_before=Pt(3))
    _add_para(ctf, "", size=6, color=GREY, space_before=Pt(4))
    _add_para(ctf, perf, size=10, color=GREEN, bold=True, space_before=Pt(2))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 – AI Dashboard Data Flow
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(sl)
_add_text(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
          "AI Detection Dashboard – Data Flow", size=28, color=ACCENT, bold=True)

ai_steps = [
    ("ai_tools table", "Stores keywords + display names\n(claude, openai, copilot, ...)\nManaged via dashboard UI", PURPLE),
    ("build_ai_s1ql()", "Generates S1QL dynamically:\nProcessName In Contains (...)\nOR ProcessCmd In Contains (...)", ORANGE),
    ("Time-Sliced Query", "8 concurrent S1 API calls\nEach slice: init → poll → events\nAll results merged + deduped", GREEN),
    ("_dedup_events()", "Groups by endpoint + process\n+ command + user\nAdds _eventCount per unique row", TEAL),
    ("Frontend Aggregation", "_aggregateAIData() computes:\nrawTotal, uniqueMatches, topApps,\ntopUsers, topEndpoints", ACCENT),
    ("Dashboard Render", "5 metric cards\nChart.js doughnut (top 25 apps)\nDense tables + drill-down overlay", RED),
]

for i, (title, desc, color) in enumerate(ai_steps):
    col = i % 3
    row = i // 3
    left = Inches(0.4 + col * 4.2)
    top = Inches(1.2 + row * 3.1)

    card = _add_box(sl, left, top, Inches(3.9), Inches(2.6))
    ctf = card.text_frame

    # number
    p0 = ctf.paragraphs[0]
    p0.text = f"  {i+1}  "
    p0.font.size = Pt(11)
    p0.font.color.rgb = BG
    p0.font.bold = True

    _add_para(ctf, title, size=13, color=color, bold=True, space_before=Pt(4))
    for line in desc.split("\n"):
        _add_para(ctf, line, size=10, color=GREY, space_before=Pt(3))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Database Schema
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(sl)
_add_text(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
          "Database Schema", size=28, color=ACCENT, bold=True)

tables = [
    ("users", "username, full_name, password_hash\nrole (admin/standard), status\nforce_password_change, last_login"),
    ("sessions", "id (random token), user_id\ncreated_at, last_active\nSliding timeout from config"),
    ("app_config", "key/value store\nS1 URL, encrypted API key\nDisk, session, password, retention"),
    ("stored_queries", "name, description, category\ndv_query (S1QL template)\nfolder_id, created_by"),
    ("query_history", "stored_query_id, query_name\nstatus, result_count, error\nuser_id, shared, executed_at"),
    ("query_results", "history_id (unique)\nresult_data (gzip blob)\nsize_bytes, created_at"),
    ("ai_tools", "keyword, display_name\ncreated_at\nDrives dynamic S1QL generation"),
    ("query_folders", "name, parent_id, sort_order\nOrganizes stored queries\ninto collapsible groups"),
]

for i, (name, desc) in enumerate(tables):
    col = i % 4
    row = i // 4
    left = Inches(0.3 + col * 3.2)
    top = Inches(1.2 + row * 3.1)

    card = _add_box(sl, left, top, Inches(3.0), Inches(2.6))
    ctf = card.text_frame
    ctf.paragraphs[0].text = name
    ctf.paragraphs[0].font.size = Pt(14)
    ctf.paragraphs[0].font.color.rgb = PURPLE
    ctf.paragraphs[0].font.bold = True
    for line in desc.split("\n"):
        _add_para(ctf, line, size=9, color=GREY, space_before=Pt(3))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Authentication Flow
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_set_bg(sl)
_add_text(sl, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
          "Authentication Flow", size=28, color=ACCENT, bold=True)

auth_steps = [
    ("1. Login", "POST /api/auth/login\nbcrypt password verify\nCheck user status = active", ACCENT),
    ("2. Session Created", "Random session token\nStored in sessions table\nSet-Cookie: sqh_session\n(httponly, samesite=lax)", GREEN),
    ("3. Every Request", "require_auth middleware\nLoads session + user\nEnforces timeout on last_active\nBumps last_active timestamp", ORANGE),
    ("4. Admin Guard", "require_admin = require_auth\n+ role == 'admin'\nProtects /api/admin/*", RED),
    ("5. Logout / Expiry", "DELETE session row\nClear cookie\nBackground cleanup every 5min", PURPLE),
]

for i, (title, desc, color) in enumerate(auth_steps):
    left = Inches(0.4 + i * 2.55)

    card = _add_box(sl, left, Inches(1.5), Inches(2.35), Inches(3.5))
    ctf = card.text_frame
    ctf.paragraphs[0].text = title
    ctf.paragraphs[0].font.size = Pt(12)
    ctf.paragraphs[0].font.color.rgb = color
    ctf.paragraphs[0].font.bold = True
    for line in desc.split("\n"):
        _add_para(ctf, line, size=9, color=GREY, space_before=Pt(4))

    if i < len(auth_steps) - 1:
        _add_text(sl, left + Inches(2.35), Inches(2.9), Inches(0.3), Inches(0.3),
                  "→", size=16, color=ARROW_CLR, bold=True)


# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
import os
out = os.path.join(os.path.dirname(__file__), "Security_Query_Hub_Data_Flow.pptx")
prs.save(out)
print(f"Saved → {out}")
